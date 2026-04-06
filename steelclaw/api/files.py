"""REST API for file uploads — handles images, documents, and audio for chat attachments."""

from __future__ import annotations

import base64
import mimetypes
import tempfile
import uuid
from pathlib import Path

from fastapi import APIRouter, File, HTTPException, Request, UploadFile

router = APIRouter()

# Supported file types and their categories
IMAGE_TYPES = {"image/jpeg", "image/png", "image/gif", "image/webp", "image/svg+xml"}
DOCUMENT_TYPES = {
    "application/pdf",
    "text/plain",
    "text/csv",
    "text/html",
    "text/markdown",
    "application/json",
    "application/xml",
    # Microsoft Office formats
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
AUDIO_TYPES = {"audio/mpeg", "audio/wav", "audio/webm", "audio/ogg", "audio/mp4", "audio/flac"}

MAX_FILE_SIZE = 20 * 1024 * 1024  # 20 MB

# Temp storage for uploaded files (keyed by file_id)
_upload_store: dict[str, dict] = {}


def _classify_mime(mime: str) -> str:
    """Classify a MIME type into a category."""
    if mime in IMAGE_TYPES:
        return "image"
    if mime in DOCUMENT_TYPES:
        return "document"
    if mime in AUDIO_TYPES:
        return "audio"
    if mime.startswith("text/"):
        return "document"
    return "unknown"


@router.post("/upload")
async def upload_file(request: Request, file: UploadFile = File(...)) -> dict:
    """Upload a file and return its ID and metadata for chat attachment."""
    content = await file.read()

    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(413, f"File too large (max {MAX_FILE_SIZE // 1024 // 1024} MB)")

    filename = file.filename or "unnamed"
    mime = file.content_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"
    category = _classify_mime(mime)

    if category == "unknown":
        raise HTTPException(
            400,
            f"Unsupported file type: {mime}. Supported: images, PDFs, text files, audio.",
        )

    file_id = str(uuid.uuid4())

    # For images: store base64 for direct LLM consumption
    # For documents: extract text content
    # For audio: transcribe to text
    processed = await _process_file(content, filename, mime, category, request)

    _upload_store[file_id] = {
        "id": file_id,
        "filename": filename,
        "mime": mime,
        "category": category,
        "size": len(content),
        "base64": base64.b64encode(content).decode() if category == "image" else None,
        "text_content": processed.get("text"),
        "preview": processed.get("preview", ""),
    }

    return {
        "id": file_id,
        "filename": filename,
        "mime": mime,
        "category": category,
        "size": len(content),
        "preview": processed.get("preview", ""),
    }


async def _process_file(
    content: bytes, filename: str, mime: str, category: str, request: Request
) -> dict:
    """Process uploaded file to extract content for the LLM."""
    if category == "image":
        return {"preview": f"[Image: {filename}]"}

    if category == "document":
        if mime == "application/pdf":
            text = _extract_pdf_text(content)
            preview = text[:200] + "..." if len(text) > 200 else text
            return {"text": text, "preview": preview}
        if mime in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or filename.lower().endswith((".docx", ".doc")):
            text = _extract_docx_text(content)
            preview = text[:200] + "..." if len(text) > 200 else text
            return {"text": text, "preview": preview}
        if mime in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ) or filename.lower().endswith((".xlsx", ".xls")):
            text = _extract_xlsx_text(content)
            preview = text[:200] + "..." if len(text) > 200 else text
            return {"text": text, "preview": preview}
        if mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or filename.lower().endswith((".pptx", ".ppt")):
            text = _extract_pptx_text(content)
            preview = text[:200] + "..." if len(text) > 200 else text
            return {"text": text, "preview": preview}
        else:
            # Text-based documents
            try:
                text = content.decode("utf-8")
            except UnicodeDecodeError:
                try:
                    text = content.decode("latin-1")
                except Exception:
                    text = "[Could not decode file content]"
            preview = text[:200] + "..." if len(text) > 200 else text
            return {"text": text, "preview": preview}

    if category == "audio":
        # Try to transcribe using the voice system
        text = await _transcribe_audio(content, filename, request)
        preview = text[:200] + "..." if len(text) > 200 else text
        return {"text": text, "preview": preview}

    return {"preview": f"[File: {filename}]"}


def _extract_docx_text(content: bytes) -> str:
    """Extract text from a DOCX file using python-docx."""
    try:
        import io
        import docx  # python-docx

        doc = docx.Document(io.BytesIO(content))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts) or "[DOCX contains no extractable text]"
    except ImportError:
        return "[DOCX reading requires python-docx — install with: pip install python-docx]"
    except Exception as e:
        return f"[Error reading DOCX: {e}]"


def _extract_xlsx_text(content: bytes) -> str:
    """Extract text from an XLSX file using openpyxl."""
    try:
        import io
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    parts.append(row_text)
        wb.close()
        return "\n".join(parts) or "[XLSX contains no extractable text]"
    except ImportError:
        return "[XLSX reading requires openpyxl — install with: pip install openpyxl]"
    except Exception as e:
        return f"[Error reading XLSX: {e}]"


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from a PPTX file using python-pptx."""
    try:
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts) or "[PPTX contains no extractable text]"
    except ImportError:
        return "[PPTX reading requires python-pptx — install with: pip install python-pptx]"
    except Exception as e:
        return f"[Error reading PPTX: {e}]"


def _extract_pdf_text(content: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=content, filetype="pdf")
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n\n".join(pages).strip() or "[PDF contains no extractable text]"
    except ImportError:
        try:
            from pypdf import PdfReader
            import io

            reader = PdfReader(io.BytesIO(content))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages).strip() or "[PDF contains no extractable text]"
        except ImportError:
            return "[PDF reading requires PyMuPDF or pypdf — install with: pip install pymupdf]"
        except Exception as e:
            return f"[Error reading PDF: {e}]"
    except Exception as e:
        return f"[Error reading PDF: {e}]"


async def _transcribe_audio(content: bytes, filename: str, request: Request) -> str:
    """Transcribe audio using the voice subsystem."""
    try:
        settings = request.app.state.settings.agents.voice
        if not settings.enabled:
            return "[Audio transcription requires voice mode to be enabled in settings]"

        from steelclaw.voice.transcription import Transcriber

        transcriber = Transcriber(settings)
        suffix = Path(filename).suffix or ".webm"
        tmp = Path(tempfile.gettempdir()) / f"sc_upload_{uuid.uuid4().hex}{suffix}"
        tmp.write_bytes(content)

        try:
            result = await transcriber.transcribe(str(tmp))
            return result.text if hasattr(result, "text") else str(result.to_dict().get("text", ""))
        finally:
            tmp.unlink(missing_ok=True)
    except Exception as e:
        return f"[Audio transcription failed: {e}]"


def get_upload(file_id: str) -> dict | None:
    """Retrieve an uploaded file's data by ID."""
    return _upload_store.get(file_id)


def cleanup_upload(file_id: str) -> None:
    """Remove an upload from the store after it's been processed."""
    _upload_store.pop(file_id, None)
